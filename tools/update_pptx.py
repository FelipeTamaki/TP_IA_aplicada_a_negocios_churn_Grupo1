"""Actualiza presentacion_tecnica_churn.pptx al cierre nuevo (Random Forest tuneado).

- Reemplaza textos (modelo final, hiperparametros, selecciones, conclusion).
- Actualiza las tablas (comparacion, sensibilidad temporal, umbral, test).
- Re-embebe las figuras regeneradas (comparacion, temporal, matriz de confusion,
  interpretabilidad combo).

Guarda en un archivo nuevo si el original esta abierto/bloqueado.
"""

from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parents[1]
SRC = ROOT / "presentacion_tecnica_churn.pptx"
FIG = ROOT / "reports" / "figures"


def set_cell(cell, value):
    para = cell.text_frame.paragraphs[0]
    if para.runs:
        para.runs[0].text = value
        for run in para.runs[1:]:
            run._r.getparent().remove(run._r)
    else:
        para.text = value


def replace_run(shape, p_idx, r_idx, new_text):
    para = shape.text_frame.paragraphs[p_idx]
    para.runs[r_idx].text = new_text


def replace_substr_in_shape(shape, old, new):
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)


def replace_pic(shape, path):
    blip = shape._element.findall(".//" + qn("a:blip"))[0]
    rid = blip.get(qn("r:embed"))
    part = shape.part.related_part(rid)
    part._blob = Path(path).read_bytes()


def shapes_by_name(slide):
    return {sh.name: sh for sh in slide.shapes}


def main():
    prs = Presentation(str(SRC))
    slides = list(prs.slides)

    # ---- Slide 1: modelo final en portada ----
    s = shapes_by_name(slides[0])
    replace_run(s["Text 12"], 0, 0, "Modelo final: Random Forest")

    # ---- Slide 7 (idx 6): hiperparametros tuneados ----
    s = shapes_by_name(slides[6])
    replace_run(s["Text 13"], 0, 0, "C = 0.1   class_weight = balanced")  # logistica
    replace_run(s["Text 17"], 1, 0, "max_depth = 6")                       # arbol
    replace_run(s["Text 21"], 0, 0, "n_estimators = 300")                  # RF
    replace_run(s["Text 21"], 1, 0, "max_depth = 12   max_features = 0.5")
    replace_run(s["Text 20"], 0, 0, "Random Forest (FINAL)")
    replace_run(
        s["Text 23"], 0, 1,
        "Los tres modelos se ajustaron con GridSearchCV sobre los mismos folds "
        "agrupados optimizando F2 (ver Modelo.ipynb, seccion 12 bis).",
    )

    # ---- Slide 9 (idx 8): sensibilidad temporal (RF tuneado) ----
    s = shapes_by_name(slides[8])
    temporal = [
        ["Completo", "0,763", "0,768", "0,743", "0,830", "0,957"],
        ["Sin Complain", "0,699", "0,703", "0,683", "0,784", "0,942"],
        ["Sin DaySinceLastOrder", "0,759", "0,765", "0,737", "0,821", "0,954"],
        ["Sin ambas (final)", "0,698", "0,707", "0,662", "0,768", "0,936"],
    ]
    tbl = next(sh.table for sh in slides[8].shapes if sh.has_table)
    for r, row_vals in enumerate(temporal, start=1):
        for c, val in enumerate(row_vals):
            set_cell(tbl.rows[r].cells[c], val)
    replace_pic(s["Image 0"], FIG / "temporal_sensitivity_f2.png")

    # ---- Slide 10 (idx 9): comparacion de modelos ----
    s = shapes_by_name(slides[9])
    replace_run(s["Text 3"], 0, 0, "Fuera de fold, modelos tuneados, umbral 0,50")
    comp = [
        ["Regresion logistica", "0,688", "0,809", "0,431", "0,609", "0,868"],
        ["Arbol de decision", "0,669", "0,780", "0,427", "0,579", "0,855"],
        ["Random Forest", "0,698", "0,707", "0,662", "0,768", "0,936"],
    ]
    tbl = next(sh.table for sh in slides[9].shapes if sh.has_table)
    for r, row_vals in enumerate(comp, start=2):
        for c, val in enumerate(row_vals):
            set_cell(tbl.rows[r].cells[c], val)
    replace_run(
        s["Text 7"], 0, 1,
        "Random Forest tuneado. Criterio primario F2 (el mas alto fuera de fold). "
        "A igual recall genera la mitad de falsas alertas que la logistica. "
        "La logistica queda como baseline interpretable.",
    )
    replace_pic(s["Image 0"], FIG / "final_model_comparison_cv.png")

    # ---- Slide 11 (idx 10): umbral y test final ----
    s = shapes_by_name(slides[10])
    tables = [sh.table for sh in slides[10].shapes if sh.has_table]
    thr_tbl, test_tbl = tables[0], tables[1]
    thr_rows = [
        ["0,27 (final)", "0,797", "0,906", "0,537", "1.280"],
        ["0,50", "0,698", "0,707", "0,662", "810"],
    ]
    for r, row_vals in enumerate(thr_rows, start=1):
        for c, val in enumerate(row_vals):
            set_cell(thr_tbl.rows[r].cells[c], val)
    test_rows = [
        ["F2", "0,817", "PR-AUC", "0,834"],
        ["Recall", "0,911", "ROC-AUC", "0,955"],
        ["Precision", "0,579", "Accuracy", "0,873"],
        ["F1", "0,708", "", ""],
    ]
    for r, row_vals in enumerate(test_rows, start=1):
        for c, val in enumerate(row_vals):
            set_cell(test_tbl.rows[r].cells[c], val)
    replace_run(s["TextBox 31"], 0, 0, "173 TP / 17 FN / 126 FP / 810 TN")
    replace_pic(s["Picture 29"], FIG / "final_test_confusion_matrix.png")

    # ---- Slide 12 (idx 11): interpretabilidad y limitaciones ----
    s = shapes_by_name(slides[11])
    for sh in slides[11].shapes:
        if sh.has_text_frame:
            replace_substr_in_shape(sh, "36,5%", "57,9%")
    replace_pic(s["Picture 7"], FIG / "final_interpretability_combo.png")

    # ---- Slide 13 (idx 12): conclusion ----
    s = shapes_by_name(slides[12])
    replace_substr_in_shape(
        s["Text 23"],
        "Modelo y umbral 0,41 fijados con F2 y recall como prioridad.",
        "Modelo y umbral 0,27 fijados con F2 y recall como prioridad.",
    )
    replace_run(
        s["Text 25"], 0, 1, "Random Forest tuneado, umbral 0,27: ",
    )
    replace_run(
        s["Text 25"], 0, 2,
        "detecta el 91,1% de los churns (recall) con F2 = 0,817 y precision 57,9% en "
        "test, reduciendo casi a la mitad las falsas alertas frente a la logistica.",
    )

    # ---- Guardar ----
    try:
        prs.save(str(SRC))
        print(f"Guardado sobre el original: {SRC.name}")
    except PermissionError:
        alt = SRC.with_name("presentacion_tecnica_churn_actualizada.pptx")
        prs.save(str(alt))
        print(f"Original bloqueado (abierto en PowerPoint). Guardado en: {alt.name}")


if __name__ == "__main__":
    main()
